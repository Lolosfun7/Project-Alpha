from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import matplotlib.pyplot as plt

# Create a presentation object
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Are You Smarter Than a 5th Grader in Stocks?"
subtitle.text = "Analyzing Knowledge of Stocks During COVID-19"

# Add data slides
data = {
    "Category": ["5th Graders", "Adults"],
    "Correct Answers (%)": [65, 85]
}

df = pd.DataFrame(data)

# Create a bar chart
plt.bar(df['Category'], df['Correct Answers (%)'], color=['blue', 'green'])
plt.title("Knowledge of Stocks During COVID-19")
plt.ylabel("Correct Answers (%)")
plt.savefig("chart.png")
plt.close()

# Add slide for graph
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Stock Knowledge Comparison"
pic = slide.shapes.add_picture("chart.png", Inches(1), Inches(1), width=Inches(8), height=Inches(4.5))

# Save the presentation
prs.save("Are_You_Smarter_Than_a_5th_Grader.pptx")
